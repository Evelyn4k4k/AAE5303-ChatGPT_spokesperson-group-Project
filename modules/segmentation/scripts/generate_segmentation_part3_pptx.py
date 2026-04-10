#!/usr/bin/env python3
"""
Build Part 3-only PowerPoint for semantic segmentation (taiwanhaitong).
Embeds real charts from figures/part3_presentation/ (run plot_segmentation_part3_figures.py first).

Output: figures/part3_presentation/AAE5303_Part3_Semantic_Segmentation.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

MODULE_ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = MODULE_ROOT / "figures" / "part3_presentation"


def _add_bullet_slide(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(17)


def _add_picture_slide(prs: Presentation, title: str, img_path: Path, left_in=0.45, top_in=1.35, width_in=12.2) -> None:
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.35), Inches(12), Inches(0.8))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(24)
    tx.text_frame.paragraphs[0].font.bold = True
    slide.shapes.add_picture(str(img_path), Inches(left_in), Inches(top_in), width=Inches(width_in))


def main() -> None:
    log_path = MODULE_ROOT / "results" / "uavscenes_training_log.json"
    metrics_path = MODULE_ROOT / "final_candidate" / "metrics_best.json"
    cfg_path = MODULE_ROOT / "configs" / "uavscenes_amtown02_interval5_14cls.yaml"

    rows = json.loads(log_path.read_text(encoding="utf-8")) if log_path.is_file() else []
    metrics = json.loads(metrics_path.read_text(encoding="utf-8")) if metrics_path.is_file() else {}
    best = max(rows, key=lambda r: float(r.get("miou", 0.0))) if rows else {}

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Section title
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "PART 03 — Semantic Segmentation"
    s.placeholders[1].text = (
        "AAE5303 — Robust Control Technology in Low-Altitude Aerial Vehicle\n"
        "GE Haitong (segmentation) · ChatGPT_spokesperson"
    )

    _add_bullet_slide(
        prs,
        "Semantic segmentation module",
        [
            "Baseline: milesial / Pytorch-UNet (U-Net, encoder–decoder + skip connections)",
            "Data: UAVScenes AMtown02, interval=5 — paired RGB + semantic masks (14 valid classes)",
            "Role in pipeline: dense per-pixel class labels for scene understanding (complements VO pose + 3D structure)",
            "Deliverables: tuned hyperparameters, checkpoints, leaderboard JSON (mIoU / Dice / fwIoU), reproducible scripts",
        ],
    )

    _add_bullet_slide(
        prs,
        "Pipeline (our implementation)",
        [
            "Download UAVScenes camera images + semantic labels → place under modules/segmentation/datasets/…",
            "analyze_mask_classes.py → remap UAVScenes sparse IDs to contiguous 0..13 (remap_uavscenes_masks.py)",
            "train_uavscenes_unet.py — AdamW, YAML config, per-epoch val metrics logged to JSON",
            "predict_unet.py — batch inference; evaluate_masks.py — full-set metrics → metrics_best.json",
            "export_leaderboard_json.py — course submission schema (percent scale)",
        ],
    )

    _add_bullet_slide(
        prs,
        "Experiment configuration (final run)",
        [
            f"Epochs: 80 | Batch size: 1 | Optimizer: AdamW | LR: 5e-5 | Weight decay: 1e-8",
            f"Input scale: 0.25 (resize) | AMP: off | Classes: 14 | Bilinear upsampling: off",
            f"Best validation mIoU (from log): {float(best.get('miou', 0)):.2f}% at epoch {int(best.get('epoch', 0))}",
            "Checkpoint for inference: selected by maximum val mIoU in uavscenes_training_log.json",
        ],
    )

    # Figures (must exist)
    for name, slide_title in [
        ("part3_loss_train_val.png", "Real data — training / validation loss"),
        ("part3_val_metrics_miou_dice_fwiou.png", "Real data — validation mIoU, Dice, fwIoU vs epoch"),
        ("part3_final_metrics_bar.png", "Real data — final test metrics (1380 pairs)"),
        ("part3_dashboard_2x2.png", "Real data — Part 3 dashboard (single slide summary)"),
    ]:
        p = FIG_DIR / name
        if p.is_file():
            _add_picture_slide(prs, slide_title, p, width_in=12.0)
        else:
            _add_bullet_slide(prs, slide_title, [f"(Run plot_segmentation_part3_figures.py — missing {name})"])

    # Same underlying JSON as part3 charts — generated earlier by plot_training_curves.py
    classic = [
        (MODULE_ROOT / "figures" / "training_loss_curve.png", "Same run — loss (plot_training_curves.py style)"),
        (MODULE_ROOT / "figures" / "validation_metrics_curve.png", "Same run — val metrics (plot_training_curves.py style)"),
    ]
    for path, slide_title in classic:
        if path.is_file():
            _add_picture_slide(prs, slide_title, path, width_in=12.0)

    _add_bullet_slide(
        prs,
        "Final results (test / leaderboard-style)",
        [
            f"mIoU: {metrics.get('miou', '—')}%",
            f"Dice: {metrics.get('dice_score', '—')}%",
            f"fwIoU: {metrics.get('fwiou', '—')}%",
            f"Evaluated pairs: {metrics.get('num_pairs', '—')} | Classes: {metrics.get('num_classes', '—')}",
            "Figures above are generated from uavscenes_training_log.json and metrics_best.json (no synthetic metrics).",
        ],
    )

    _add_bullet_slide(
        prs,
        "Key findings",
        [
            "Label remapping (sparse UAVScenes IDs → 0..13) is required for stable multi-class training.",
            "AdamW + low LR + small input scale matches course UNet demo recommendations for this sequence.",
            "Val mIoU peaks before the last epoch — model selection by best val checkpoint matters for reporting.",
            "Fair comparison requires the same mask protocol and the official UAVScenes interval=5 split size.",
        ],
    )

    out = FIG_DIR / "AAE5303_Part3_Semantic_Segmentation.pptx"
    FIG_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
