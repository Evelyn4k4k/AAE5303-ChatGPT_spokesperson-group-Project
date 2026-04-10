"""
Generate a technical overview PowerPoint for AAE5303 group project presentation.
Requires: pip install python-pptx

Output: modules/segmentation/figures/AAE5303_Group_Project_Technical_Overview.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def _bullet(slide, title: str, lines: list[str]) -> None:
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def main() -> None:
    root = Path(__file__).resolve().parents[1]
    out = root / "figures" / "AAE5303_Group_Project_Technical_Overview.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Title ---
    slide_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(slide_layout)
    s.shapes.title.text = "AAE5303 Group Project"
    s.placeholders[1].text = (
        "Technical stack: VO + 3D reconstruction + semantic segmentation\n"
        "Course: Robust Control Technology in Low-Altitude Aerial Vehicle\n"
        "Repository: ChatGPT_spokesperson group project"
    )

    # --- Course framing ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Course requirements (grading context)",
        [
            "Hyperparameter tuning across three tracks: Visual Odometry, 3D reconstruction, semantic segmentation",
            "Leaderboard submission (reproducible metrics) + group/individual presentation",
            "Private GitHub repo with Lead TA & instructor added as collaborators for reproduction checks",
            "Suggested baselines: ORB-SLAM3 (VO), OpenSplat (3D), Pytorch-UNet (segmentation)",
        ],
    )

    # --- Three modules ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Repository modular structure",
        [
            "modules/vo — Visual odometry (Evelyn4k4k)",
            "modules/reconstruction — 3D scene reconstruction (wymmust)",
            "modules/segmentation — Semantic segmentation (taiwanhaitong)",
            "Shared sequence alignment: AMtown02 (MARS-LVIG / UAVScenes family) where applicable",
        ],
    )

    # --- Dataset ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Datasets & alignment",
        [
            "MARS-LVIG / UAVScenes ecosystem for low-altitude aerial perception",
            "VO: monocular RGB on AMtown02; trajectory vs ground truth",
            "Segmentation evaluation: UAVScenes AMtown02, interval=5, paired RGB + semantic masks",
            "Interval=5: subsampled frames for manageable training/eval size while keeping scene coverage",
        ],
    )

    # --- VO ORB-SLAM3 ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Module 1 — Visual odometry (ORB-SLAM3)",
        [
            "Backend: ORB-SLAM3 — feature-based monocular SLAM / VO",
            "ORB features: fast binary descriptors, rotation-aware, suitable for real-time tracking",
            "Pipeline: image extraction → camera/ORB tuning → run SLAM → export CameraTrajectory / KeyFrameTrajectory",
            "Hyperparameters: camera intrinsics (official AMtown calibration), ORB density (e.g. medium preset), resolution",
            "Evaluation: ATE RMSE, RPE translation/rotation drift, trajectory completeness vs GT timestamps",
        ],
    )

    # --- 3D OpenSplat ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Module 2 — 3D reconstruction (course baseline: OpenSplat)",
        [
            "OpenSplat: open-source 3D reconstruction aligned with modern Gaussian splatting workflows",
            "Goal: recover dense / viewable 3D representation of the scene from images or posed views",
            "Typical knobs for tuning: sampling density, optimization iterations, initialization, regularization",
            "Outputs: 3D model / splat export for viewer (e.g. web Model Viewer) — details per group implementation",
            "Coordinate with VO outputs (poses) or image sets from the same AMtown02 extract where possible",
        ],
    )

    # --- Segmentation intro ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Module 3 — Semantic segmentation (Pytorch-UNet)",
        [
            "Task: per-pixel class label for each aerial image (multi-class semantic segmentation)",
            "Baseline implementation: milesial / Pytorch-UNet (U-Net in PyTorch)",
            "U-Net: encoder–decoder with skip connections; strong baseline for dense prediction",
            "Training target: UAVScenes-style labels; course example uses 14 valid semantic classes on AMtown02 interval=5",
        ],
    )

    # --- U-Net detail ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Segmentation — model & optimization",
        [
            "Network: U-Net — multi-scale feature fusion for precise boundaries",
            "Framework: PyTorch; GPU training with CUDA; optional mixed precision (AMP) for speed/memory",
            "Optimizer: AdamW (decoupled weight decay) — common for vision transformers & modern CNN training",
            "Hyperparameters tuned: learning rate, batch size, epochs, input scale (resize factor), AMP on/off",
            "Checkpointing & resume: long runs saved per epoch; best model selected by validation mIoU",
        ],
    )

    # --- Data preprocessing ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Segmentation — data & label mapping",
        [
            "UAVScenes masks use dataset-specific semantic IDs (sparse / non-contiguous)",
            "Remapping: map chosen valid class IDs to contiguous 0..C−1 (here C=14) for softmax training",
            "Scripts: analyze unique IDs in mask folder → deterministic remap → paired img/mask filenames",
            "Optional bootstrap: pseudo-labels from RGB heuristics (not for final leaderboard claims)",
        ],
    )

    # --- Metrics ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Segmentation — evaluation metrics (leaderboard)",
        [
            "mIoU (mean Intersection-over-Union): class-averaged overlap — primary semantic segmentation metric",
            "Dice / F1-style overlap: agreement between prediction and GT regions",
            "fwIoU (frequency-weighted IoU): IoU weighted by class pixel frequency (handles imbalance)",
            "All reported in percent (%) to match course leaderboard JSON schema",
        ],
    )

    # --- Pipeline automation ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Segmentation — engineering & reproducibility",
        [
            "YAML configs for experiment settings; CSV template for hyperparameter sweep log",
            "Training log JSON (per-epoch loss & val metrics); matplotlib curves for presentation",
            "Inference: sliding-window / full-image prediction → PNG masks; batch prediction over eval set",
            "export_leaderboard_json.py: metrics_best.json → ChatGPT_spokesperson_leaderboard.json",
        ],
    )

    # --- Extra tools ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Extra / supporting technology",
        [
            "Python virtual environment (venv) — isolated dependencies",
            "Git / GitHub — version control, branch per contributor, private repo for grading",
            "tqdm — live training progress; Pillow — image I/O",
            "NumPy — metric aggregation; PyYAML — configuration files",
            "Optional: pseudo-label pipeline, smoke-test synthetic data for CI-style sanity checks only",
        ],
    )

    # --- Example results placeholder ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Segmentation — example final numbers (replace if your run differs)",
        [
            "1380 image–mask pairs evaluated, 14 classes",
            "mIoU ≈ 78.24% | Dice ≈ 87.28% | fwIoU ≈ 89.65%",
            "Figures: training_loss_curve.png, validation_metrics_curve.png under modules/segmentation/figures/",
            "Insert qualitative overlays (RGB + colored mask) on this slide or next for visual proof",
        ],
    )

    # --- Summary ---
    s = prs.slides.add_slide(prs.slide_layouts[1])
    _bullet(
        s,
        "Summary",
        [
            "Three open baselines cover perception stack: pose (ORB-SLAM3), 3D (OpenSplat), semantics (U-Net)",
            "Common theme: hyperparameter tuning + reproducible scripts + leaderboard-compatible metrics",
            "References: ORB-SLAM3 repo, OpenSplat, milesial/Pytorch-UNet, UAVScenes, MARS-LVIG, course leaderboard demo repos",
        ],
    )

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"Saved: {out}")


if __name__ == "__main__":
    main()
